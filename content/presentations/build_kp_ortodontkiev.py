# -*- coding: utf-8 -*-
"""Комерційна пропозиція (PPC · Google Ads) для ortodontkiev.com.ua, UA, 11 слайдів.

Стиль і дизайн-токени успадковані з build_site24_short_deck.py (офіційний бренд-скрипт
Site24) та content/presentations/2026-06-25-EN-site24-presentation-guidelines.html
(секція «08 · Шаблони · комерційна пропозиція / пітч»): Ubuntu, --ink/--blue/--green/
--mint/--tint, макет Титул → Стратегія → Кейс → Кейси → Доказ → План → Тарифи → CTA.

Дані:
- План робіт і трудовитрати — з таблиці клієнтки (Google Sheets, план для ortodontkiev.com.ua).
- Бюджети на кліки (грн/день, aver. CPC) — надані Катериною для 4 кампаній.
- Кейси й «коротко про нас» — з еталонного деку (Google Slides, PPC КП іншого клієнта) та
  канонічного гросарію content/translations/00-glossary.md (9 років, 1000+ клієнтів,
  86% успішних проєктів, 30+ ніш) — уникаємо неперевірених цифр за правилами CLAUDE.md.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(os.path.dirname(BASE))
BADGE = os.path.join(ROOT, "assets", "badges")
NICHE = os.path.join(ROOT, "assets", "niches")

INK = RGBColor(0x0E, 0x0E, 0x3A); BLUE = RGBColor(0x29, 0x66, 0xFF)
GREEN = RGBColor(0x36, 0xEF, 0x74); GREENM = RGBColor(0x94, 0xFC, 0xB1)
MINT = RGBColor(0xE6, 0xFA, 0xEE); TINT = RGBColor(0xDA, 0xF8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); INK60 = RGBColor(0x6E, 0x6E, 0x88)
LINE = RGBColor(0xDD, 0xE1, 0xEC); INK12 = RGBColor(0xCE, 0xD2, 0xDE)
FONT = "Ubuntu"

prs = Presentation(); prs.slide_width = Pt(960); prs.slide_height = Pt(540)
BLANK = prs.slide_layouts[6]
MX = 64


def imsize(p):
    with Image.open(p) as im:
        return im.size


def fit_pic(s, path, bx, by, bw, bh):
    iw, ih = imsize(path); r = min(bw / iw, bh / ih); w = iw * r; h = ih * r
    s.shapes.add_picture(path, Pt(bx + (bw - w) / 2), Pt(by + (bh - h) / 2), width=Pt(w), height=Pt(h))


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill=WHITE, line=None, radius=0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = min(0.5, radius / float(min(w, h)))
        except Exception:
            pass
    shp.shadow.inherit = False
    return shp


def txt(s, x, y, w, h, runs, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.15, italic=False):
    tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h)); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0; tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = ls
    for run in runs:
        t, c, b = run[0], run[1], run[2]
        it = run[3] if len(run) > 3 else italic
        r = p.add_run(); r.text = t; r.font.size = Pt(size); r.font.bold = b; r.font.italic = it
        r.font.name = FONT; r.font.color.rgb = c
    return tb


def para(s, x, y, w, h, parts, size=16, color=INK, ls=1.4, anchor=MSO_ANCHOR.MIDDLE):
    return txt(s, x, y, w, h, parts, size=size, color=color, ls=ls, anchor=anchor)


def heading(s, x, y, w, parts, size=30, h=48):
    return txt(s, x, y, w, h, parts, size=size, bold=True, ls=1.08)


def pic_h(s, path, x, y, h):
    iw, ih = imsize(path); s.shapes.add_picture(path, Pt(x), Pt(y), height=Pt(h)); return h * iw / ih


def pill(s, x, y, text, size=13, fill=BLUE, color=WHITE):
    tw = len(text) * size * 0.62 + 30; hh = size + 14
    rect(s, x, y, tw, hh, fill=fill, radius=hh / 2)
    txt(s, x, y, tw, hh, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return tw, hh


def footer(s, num):
    txt(s, MX, 504, 500, 20, "site24.com.ua", size=12, color=INK60, align=PP_ALIGN.LEFT)
    txt(s, 760, 504, 136, 20, str(num), size=12, color=INK60, align=PP_ALIGN.RIGHT)


def link(s, x, y, w, text, url, size=13, color=BLUE, align=PP_ALIGN.LEFT, bold=True):
    tb = s.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(22)); tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = color
    r.hyperlink.address = url
    return tb


def numbered_item(s, x, y, w, n, text, size=13.5, ls=1.32, box=30):
    """Мала tint-плашка з синім номером + текст пункту плану."""
    rect(s, x, y, box, box, fill=TINT, radius=8)
    txt(s, x, y, box, box, str(n), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + box + 18, y - 3, w - box - 18, 90, text, size=size, color=INK60, ls=ls, anchor=MSO_ANCHOR.TOP)


def case_slide(tagtxt, htitle, subtitle, stats4, task, case_url, foot_num):
    s = slide(WHITE)
    pill(s, MX, 50, tagtxt, size=13)
    heading(s, MX, 80, 832, htitle, size=27, h=50)
    para(s, MX, 134, 832, 42, [(subtitle, INK60, False)], size=14, anchor=MSO_ANCHOR.TOP, ls=1.3)
    rect(s, MX, 182, 832, 60, fill=MINT, radius=14)
    rect(s, MX + 22, 202, 4, 22, fill=BLUE)
    para(s, MX + 40, 182, 772, 60, [("Задача: ", INK, True), (task, INK60, False)], size=13, anchor=MSO_ANCHOR.MIDDLE, ls=1.24)
    gx, gy = MX, 258; cw, ch = 200, 140
    for i, (n, l) in enumerate(stats4):
        x = gx + i * (cw + 11)
        rect(s, x, gy, cw, ch, fill=TINT, radius=18)
        rect(s, x + 22, gy + 20, 14, 14, fill=GREEN)
        txt(s, x + 22, gy + 40, cw - 40, 48, n, size=32, bold=True, color=BLUE, anchor=MSO_ANCHOR.TOP)
        txt(s, x + 22, gy + 90, cw - 36, 48, l, size=11.5, color=INK60, ls=1.18)
    link(s, MX, 424, 320, "Детальніше про кейс  →", case_url, size=14, color=BLUE)
    link(s, 512, 424, 320, "Читати всі кейси  →", "https://site24.com.ua/tag/kejsi/", size=14, color=INK60, align=PP_ALIGN.RIGHT)
    footer(s, foot_num)


# ===================================================================== 01 ТИТУЛ
s = slide(TINT)
pill(s, MX, 56, "PPC-РЕКЛАМА · GOOGLE ADS", size=13)
heading(s, MX, 96, 830, [("Комерційна пропозиція для ", INK, True), ("ortodontkiev.com.ua", BLUE, True)], size=38, h=110)
para(s, MX, 210, 760, 76,
     [("Ціль: лідогенерація через Google Ads на послуги ортодонтії — консультацію ортодонта, "
       "брекети, елайнери та ортодонтичні апарати.", INK60, False)],
     size=16, anchor=MSO_ANCHOR.TOP, ls=1.4)
rect(s, MX, 300, 300, 4, fill=GREEN)
chips = [("$504", "налаштування · 1-й місяць"), ("$324", "оптимізація · щомісяця"), ("22 500 ₴", "бюджет на кліки / міс")]
gx, gy = MX, 330; cw, ch = 260, 118
for i, (n, l) in enumerate(chips):
    x = gx + i * (cw + 16)
    rect(s, x, gy, cw, ch, fill=WHITE, radius=16)
    txt(s, x + 22, gy + 16, cw - 40, 44, n, size=27, bold=True, color=BLUE, anchor=MSO_ANCHOR.TOP)
    txt(s, x + 22, gy + 64, cw - 40, 44, l, size=12, color=INK60, ls=1.2, anchor=MSO_ANCHOR.TOP)
txt(s, MX, 504, 500, 20, [("серпень 2026", INK60, True), ("      ·      ", INK12, False), ("site24.com.ua", INK60, False)], size=13)

# ===================================================================== 02 СТРАТЕГІЯ РОБОТИ
s = slide(WHITE)
heading(s, MX, 50, 860, [("Стратегія ", INK, True), ("роботи", BLUE, True)], size=30)
strat = [
    ("Пошукові оголошення",
     "Налаштовуємо рекламні пошукові оголошення за релевантними запитами користувачів: "
     "«брекети Київ», «елайнери ціна», «консультація ортодонта» тощо."),
    ("Акцент на перевагах клініки",
     "В оголошеннях і на посадкових сторінках виділяємо переваги ortodontkiev.com.ua, "
     "щоб вигідно відрізнятися від конкурентів у видачі."),
    ("Розширення оголошень",
     "Використовуємо розширення (уточнення, номер телефону, адреса, посилання) для "
     "більшої привабливості в пошуку та відповідності запитам."),
]
gx, gy = MX, 122; cw, ch = 270, 252
for i, (t, d) in enumerate(strat):
    x = gx + i * (cw + 11)
    rect(s, x, gy, cw, ch, fill=MINT, radius=18)
    rect(s, x + 20, gy + 22, 38, 4, fill=BLUE)
    txt(s, x + 20, gy + 42, cw - 38, 60, t, size=17, bold=True, color=INK, ls=1.14)
    txt(s, x + 20, gy + 108, cw - 38, ch - 122, d, size=12.5, color=INK60, ls=1.35)
footer(s, "02")

# ===================================================================== 03 PPC-КЕЙС (стоматологія, Нью-Йорк)
case_slide(
    "PPC-КЕЙС · РЕЛЕВАНТНИЙ ДОСВІД",
    [("Кейс: ", INK, True), ("стоматологічна клініка в Нью-Йорку", BLUE, True)],
    "Комплексна стоматологія (імпланти, вініри, невідкладна допомога), Квінз, Нью-Йорк. "
    "Пошукова реклама Google на висококонкурентний ринок США.",
    [("636", "заявок на місяць — з 5 на старті"),
     ("13,26%", "CTR пошукових оголошень"),
     ("$25", "ціна конверсії — було $574 на старті"),
     ("1 975", "конверсій за 6 місяців")],
    "залучити максимум нових пацієнтів через пошукову рекламу на висококонкурентному ринку — "
    "медична тематика з тими самими викликами: висока вартість ліда, довгий цикл рішення, довіра до лікаря.",
    "https://site24.com.ua/ppc-reklama-stomatologichnoi-kliniki-na-rinok-ssha/",
    "03")

# ===================================================================== 04 ЩЕ НАШІ PPC-КЕЙСИ
s = slide(WHITE)
heading(s, MX, 50, 860, [("Ще наші ", INK, True), ("PPC-кейси", BLUE, True)], size=30)
cases = [
    ("Studix.eu", "Освіта за кордоном — набір у польські ВНЗ",
     [("+251%", "конверсій (45 → 158)"), ("−39%", "ціна конверсії (268 → 164 ₴)"),
      ("158", "лідів за 4 місяці"), ("4 міс", "тривалість проєкту")],
     "https://site24.com.ua/yak-zbilshyty-kilkist-konversij-iz-kontekstnoyi-reklamy-na-251-ta-poproshhatysya-z-kliyentom-kejs-studix-eu/"),
    ("The Tea", "E-commerce — китайський чай (thetea.ua)",
     [("×2,5", "дохід за рік"), ("×2,3", "транзакції (940 → 2143)"),
      ("+200%", "ROAS (333% → 523%)"), ("30+", "нових кампаній Google Ads")],
     "https://site24.com.ua/ppc-dlya-e-commerce-zrostannya-dohodu-v-25-razi-za-odin-rik/"),
]
gx, gy = MX, 118; cw, ch = 406, 300
for i, (brand, desc, stats, url) in enumerate(cases):
    x = gx + i * (cw + 20)
    rect(s, x, gy, cw, ch, fill=MINT, radius=18)
    txt(s, x + 26, gy + 22, cw - 52, 32, brand, size=19, bold=True, color=INK)
    txt(s, x + 26, gy + 58, cw - 52, 34, desc, size=12.5, color=INK60, ls=1.2)
    sx, sy = x + 26, gy + 104; scw, sch = (cw - 52 - 14) / 2, 86
    for j, (n, l) in enumerate(stats):
        sx2 = sx + (j % 2) * (scw + 14); sy2 = sy + (j // 2) * (sch + 10)
        txt(s, sx2, sy2, scw, 32, n, size=21, bold=True, color=BLUE)
        txt(s, sx2, sy2 + 34, scw, 40, l, size=10.5, color=INK60, ls=1.16)
    link(s, x + 26, gy + ch - 40, cw - 52, "Детальніше про кейс  →", url, size=12.5, color=BLUE)
link(s, MX, 434, 400, "Читати всі кейси на сайті  →", "https://site24.com.ua/tag/kejsi/", size=13, color=INK60)
footer(s, "04")

# ===================================================================== 05 КОРОТКО ПРО НАС
s = slide(TINT)
heading(s, MX, 50, 860, [("Контекстна реклама від ", INK, True), ("Site24", BLUE, True), (" у цифрах", INK, True)], size=28)
stats = [("9", "років на ринку"), ("1000+", "клієнтів"), ("86%", "успішних проєктів"), ("30+", "ніш бізнесу")]
gx, gy = MX, 130; cw, ch = 200, 150
for i, (n, l) in enumerate(stats):
    x = gx + i * (cw + 11)
    rect(s, x, gy, cw, ch, fill=WHITE, radius=18)
    rect(s, x + 22, gy + 26, 14, 14, fill=GREEN)
    txt(s, x + 22, gy + 46, cw - 40, 54, n, size=42, bold=True, color=BLUE, anchor=MSO_ANCHOR.TOP)
    txt(s, x + 22, gy + 100, cw - 36, 48, l, size=12.5, color=INK60, ls=1.18)
rect(s, MX, 316, 832, 96, fill=BLUE, radius=16)
fit_pic(s, os.path.join(NICHE, "medicine.png"), MX + 22, 316 + 24, 48, 48)
para(s, MX + 92, 316, 590, 96,
     [("Медична / стоматологічна тематика — одна з наших профільних ніш: ", WHITE, False),
      ("розуміємо специфіку довіри пацієнта, вартості ліда та роботи з call-tracking.", GREENM, False)],
     14.5, anchor=MSO_ANCHOR.MIDDLE, ls=1.35)
rect(s, 760, 330, 108, 68, fill=WHITE, radius=12)
fit_pic(s, os.path.join(BADGE, "google-partner.png"), 770, 340, 88, 48)
footer(s, "05")

# ===================================================================== 06 ПЛАН РОБОТИ. НАЛАШТУВАННЯ — 1 МІСЯЦЬ
s = slide(WHITE)
heading(s, MX, 50, 860, [("План роботи. ", INK, True), ("Налаштування реклами — 1 місяць", BLUE, True)], size=25)
setup_items = [
    "Створення акаунтів GA4 (Google Analytics), Google Ads, Google Tag Manager (GTM).",
    "Налаштування аналітики та відстеження конверсій: цілі в GA4, зв'язка Google Ads і Analytics, "
    "передача конверсійних дій.",
    "Звітність, комунікація з акаунт-менеджером — увесь місяць.",
    "Створення категорійних пошукових кампаній «Консультація ортодонта», «Брекети», «Елайнери», "
    "«Ортодонтичні апарати»: збір ключів, розширення семантики, мінус-слова, кластеризація груп, "
    "тексти оголошень (15 заголовків, 5 описів, URL, лого, зображення), стратегія ставок, модерація.",
    "Аналіз роботи акаунту: покази, аукціони, конкуренти, точки зростання — раз на місяць.",
]
y = 110
for i, item in enumerate(setup_items, 1):
    h = 58 if i != 4 else 92
    numbered_item(s, MX, y, 832, i, item, box=30)
    y += h
rect(s, MX, 462, 832, 46, fill=MINT, radius=10)
para(s, MX + 24, 462, 800, 46, [("Разом: 28 год", INK, True), ("    ·    Вартість: $504 · одноразово", INK60, False)], size=13)
footer(s, "06")

# ===================================================================== 07 ПЛАН РОБОТИ. ЩОМІСЯЧНА ОПТИМІЗАЦІЯ
s = slide(WHITE)
heading(s, MX, 50, 860, [("План роботи. ", INK, True), ("Щомісячна оптимізація", BLUE, True)], size=27)
opt_items = [
    "Звітність, комунікація з акаунт-менеджером — увесь місяць.",
    "Ведення пошукової реклами (Search + Brand): оптимізація ключових слів, аналіз пошукових запитів "
    "і мінус-слова, оптимізація текстів оголошень, тестування автостратегій, перерозподіл бюджетів — "
    "раз на тиждень; перевірка форм на сайті — раз на два тижні; контроль балансу кабінету — щоденно.",
    "Аналіз роботи облікового запису: покази, аукціони, конкуренти, точки зростання — раз на місяць.",
]
y = 118
heights = [58, 108, 58]
for i, (item, h) in enumerate(zip(opt_items, heights), 1):
    numbered_item(s, MX, y, 832, i, item, box=30)
    y += h + 18
rect(s, MX, 456, 832, 46, fill=MINT, radius=10)
para(s, MX + 24, 456, 800, 46, [("Разом: 18 год / міс", INK, True), ("    ·    Бюджет: $324 / місяць", INK60, False)], size=13)
footer(s, "07")

# ===================================================================== 08 БЮДЖЕТ НА РОБОТИ (тарифи)
s = slide(WHITE)
heading(s, MX, 50, 860, [("Бюджет на ", INK, True), ("роботи", BLUE, True)], size=30)
stages = [("1 ЕТАП", "Налаштування рекламних кампаній", "$504", "1 місяць · одноразово"),
          ("2 ЕТАП", "Щомісячна оптимізація реклами", "$324", "щомісячно")]
gx, gy = MX, 130; cw, ch = 406, 300
for i, (tag, title, price, note) in enumerate(stages):
    x = gx + i * (cw + 20)
    rect(s, x, gy, cw, ch, fill=WHITE, line=LINE, radius=18)
    rect(s, x, gy, cw, 6, fill=BLUE)
    txt(s, x + 28, gy + 34, cw - 56, 22, tag, size=13, bold=True, color=BLUE)
    txt(s, x + 28, gy + 62, cw - 56, 70, title, size=19, bold=True, color=INK, ls=1.15)
    txt(s, x + 28, gy + 152, cw - 56, 66, price, size=42, bold=True, color=INK)
    txt(s, x + 28, gy + 226, cw - 56, 30, note, size=13, color=INK60)
para(s, MX, 452, 832, 40,
     [("* Вказана вартість послуг Site24. Бюджет на рекламні кліки в Google Ads оплачується "
       "клієнтом окремо — див. наступний слайд.", INK60, False)], size=11.5, anchor=MSO_ANCHOR.TOP, ls=1.3)
footer(s, "08")

# ===================================================================== 09 ПРОРАХУНОК БЮДЖЕТУ НА КЛІКИ
s = slide(WHITE)
heading(s, MX, 50, 860, [("Прорахунок бюджету на ", INK, True), ("кліки", BLUE, True)], size=27)
head = ["Кампанія", "Бюджет / день, ₴", "Aver. CPC, ₴", "Кліків / міс", "Бюджет / міс, ₴"]
rows = [
    ("Ортодонт", "150", "10", "450", "4 500"),
    ("Брекети", "200", "13,3", "451", "6 000"),
    ("Елайнери", "250", "22", "341", "7 500"),
    ("Ортодонтичні апарати", "150", "4", "1 125", "4 500"),
]
totals = ("Разом", "750", "—", "2 367", "22 500")
colw = [258, 144, 144, 143, 143]
gx, gy = MX, 118; rh = 42
x = gx
for i, htxt in enumerate(head):
    rect(s, x, gy, colw[i], rh, fill=INK)
    align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
    txt(s, x + 14, gy, colw[i] - 24, rh, htxt, size=12, bold=True, color=WHITE, align=align, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[i]
y = gy + rh
for ri, row in enumerate(rows):
    bg = WHITE if ri % 2 == 0 else MINT
    x = gx
    for ci, cell in enumerate(row):
        rect(s, x, y, colw[ci], rh, fill=bg, line=LINE)
        align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        txt(s, x + 14, y, colw[ci] - 24, rh, cell, size=13, bold=(ci == 0), color=INK, align=align, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[ci]
    y += rh
x = gx
for ci, cell in enumerate(totals):
    rect(s, x, y, colw[ci], rh, fill=GREEN)
    align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
    txt(s, x + 14, y, colw[ci] - 24, rh, cell, size=13, bold=True, color=INK, align=align, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[ci]
y += rh + 22
para(s, MX, y, 832, 50,
     [("* Бюджет на кліки — окремо від вартості послуг Site24, оплачується напряму в Google Ads. "
       "Середня ціна за клік (Aver. CPC) орієнтовна та залежить від аукціону; кількість кліків "
       "розрахована виходячи з денного бюджету напряму.", INK60, False)], size=11.5, anchor=MSO_ANCHOR.TOP, ls=1.32)
footer(s, "09")

# ===================================================================== 10 НАШІ КЛІЄНТИ
s = slide(WHITE)
heading(s, MX, 54, 860, [("Нам ", INK, True), ("довіряють", BLUE, True)], size=30)
para(s, MX, 110, 832, 26, [("Ми ведемо контекстну рекламу для бізнесів у різних нішах — від медицини й освіти до e-commerce та послуг.", INK60, False)],
     size=15, anchor=MSO_ANCHOR.TOP)
clients = ["Ortodontkiev.com.ua", "Studix.eu", "The Tea (thetea.ua)", "Стоматологія · Нью-Йорк, США", "1000+ клієнтів у 30+ нішах"]
gx, gy = MX, 170; cw, ch = 270, 96
for i, c in enumerate(clients):
    x = gx + (i % 3) * (cw + 11); y = gy + (i // 3) * (ch + 14)
    rect(s, x, y, cw, ch, fill=MINT, radius=14)
    txt(s, x + 22, y, cw - 44, ch, c, size=14.5, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE, ls=1.2)
footer(s, "10")

# ===================================================================== 11 КОНТАКТИ
s = slide(TINT)
heading(s, MX, 130, 760, [("Готові запустити рекламу для ", INK, True), ("ortodontkiev.com.ua", BLUE, True), ("?", INK, True)], size=32, h=100)
para(s, MX, 232, 700, 30, [("Погодимо стратегію та запустимо перші кампанії протягом 3–5 днів після старту.", INK60, False)],
     size=16, anchor=MSO_ANCHOR.TOP)


def lab(x, y, t):
    txt(s, x, y, 360, 18, t, size=12.5, bold=True, color=BLUE)


def val(x, y, parts, sz=18):
    txt(s, x, y, 420, 26, parts, size=sz, color=INK)


lab(MX, 300, "КЕРІВНИК АГЕНЦІЇ")
val(MX, 322, "Катерина Золотарьова")
lab(MX, 376, "ТЕЛЕФОН")
val(MX, 398, "+38 (098) 738 77 08")
lab(500, 376, "САЙТ")
val(500, 398, [("site24.com.ua", BLUE, True)])
TG = "https://t.me/katerinazolotaryova"
btn = rect(s, 500, 300, 300, 52, fill=GREEN, radius=12)
btn.click_action.hyperlink.address = TG
btxt = txt(s, 500, 300, 268, 52, "Замовити консультацію", size=16, bold=True, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
btxt.text_frame.paragraphs[0].runs[0].hyperlink.address = TG
txt(s, 750, 316, 26, 22, "→", size=18, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
txt(s, MX, 504, 500, 20, "site24.com.ua", size=12, color=INK60)
txt(s, 760, 504, 136, 20, "11", size=12, color=INK60, align=PP_ALIGN.RIGHT)

out = os.path.join(BASE, "2026-08-25-kp-ortodontkiev.pptx")
prs.save(out)
print("SAVED", out, "| slides", len(prs.slides._sldIdLst))
