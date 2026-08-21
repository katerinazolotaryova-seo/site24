# -*- coding: utf-8 -*-
"""Коротка презентація «Site24 — про компанію» (стиль сайту site24.com.ua), UA, 9 слайдів.
Формат: як у Bettertone-деку — редаговані текстові поля/фігури, 960x540 (1px=1pt), Ubuntu.
Дизайн-токени й хелпери успадковані з build_bettertone_deck.py + гайдлайну презентацій.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
AST = os.path.join(BASE, "bettertone-assets"); IC = os.path.join(AST, "icons")
S24 = os.path.join(BASE, "site24-assets"); BADGE=os.path.join(S24,"badges"); NICHE=os.path.join(S24,"niches")

def fit_pic(s,path,bx,by,bw,bh):
    """Вписати картинку в бокс (contain), відцентрувати."""
    iw,ih=imsize(path); r=min(bw/iw,bh/ih); w=iw*r; h=ih*r
    s.shapes.add_picture(path,Pt(bx+(bw-w)/2),Pt(by+(bh-h)/2),width=Pt(w),height=Pt(h))

INK   = RGBColor(0x0E,0x0E,0x3A); BLUE = RGBColor(0x29,0x66,0xFF)
GREEN = RGBColor(0x36,0xEF,0x74); GREENM=RGBColor(0x94,0xFC,0xB1)
MINT  = RGBColor(0xE6,0xFA,0xEE); TINT = RGBColor(0xDA,0xF8,0xFF)
WHITE = RGBColor(0xFF,0xFF,0xFF); INK60= RGBColor(0x6E,0x6E,0x88)
LINE  = RGBColor(0xDD,0xE1,0xEC); INK12=RGBColor(0xCE,0xD2,0xDE)
FONT="Ubuntu"

prs=Presentation(); prs.slide_width=Pt(960); prs.slide_height=Pt(540)
BLANK=prs.slide_layouts[6]

def imsize(p):
    with Image.open(p) as im: return im.size
def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=bg; return s
def rect(s,x,y,w,h,fill=WHITE,line=None,radius=0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Pt(x),Pt(y),Pt(w),Pt(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1)
    if radius:
        try: shp.adjustments[0]=min(0.5,radius/float(min(w,h)))
        except Exception: pass
    shp.shadow.inherit=False; return shp
def txt(s,x,y,w,h,runs,size=16,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,ls=1.15):
    tb=s.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=anchor
    if isinstance(runs,str): runs=[(runs,color,bold)]
    p=tf.paragraphs[0]; p.alignment=align; p.line_spacing=ls
    for t,c,b in runs:
        r=p.add_run(); r.text=t; r.font.size=Pt(size); r.font.bold=b; r.font.name=FONT; r.font.color.rgb=c
    return tb
def para(s,x,y,w,h,parts,size=16,color=INK,ls=1.4,anchor=MSO_ANCHOR.MIDDLE):
    return txt(s,x,y,w,h,parts,size=size,color=color,ls=ls,anchor=anchor)
def bullets(s,x,y,w,items,size=17,color=INK,gap=10,ls=1.2,arrow=BLUE):
    tb=s.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(540)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=ls; p.space_after=Pt(gap)
        r1=p.add_run(); r1.text="→  "; r1.font.size=Pt(size); r1.font.bold=True; r1.font.name=FONT; r1.font.color.rgb=arrow
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(size); r2.font.name=FONT; r2.font.color.rgb=color
    return tb
def heading(s,x,y,w,parts,size=30,h=48): return txt(s,x,y,w,h,parts,size=size,bold=True,ls=1.08)
def pic_h(s,path,x,y,h):
    iw,ih=imsize(path); s.shapes.add_picture(path,Pt(x),Pt(y),height=Pt(h)); return h*iw/ih
def pill(s,x,y,text,size=13,fill=BLUE,color=WHITE):
    tw=len(text)*size*0.62+30; hh=size+14
    rect(s,x,y,tw,hh,fill=fill,radius=hh/2)
    txt(s,x,y,tw,hh,text,size=size,bold=True,color=color,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    return tw,hh
def footer(s,num):
    txt(s,60,504,140,20,str(num),size=12,color=INK60,align=PP_ALIGN.LEFT)
    logop=os.path.join(IC,"site24-ink.png")
    iw,ih=imsize(logop); w=20*iw/ih
    s.shapes.add_picture(logop,Pt(900-w),Pt(505),height=Pt(20))
def link(s,x,y,w,text,url,size=13,color=BLUE,align=PP_ALIGN.LEFT,bold=True):
    """Клікабельний текст-лінк."""
    tb=s.shapes.add_textbox(Pt(x),Pt(y),Pt(w),Pt(22)); tf=tb.text_frame; tf.word_wrap=False
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.name=FONT; r.font.color.rgb=color
    r.hyperlink.address=url
    return tb
MX=64

def case_slide(tagtxt,htitle,subtitle,stats4,task,case_url,foot_num):
    """Слайд-кейс: тип-піл, заголовок, ніша, 4 стат-картки, банер-задача, лінки."""
    s=slide(WHITE)
    pill(s,MX,50,tagtxt,size=13)
    heading(s,MX,80,832,htitle,size=27,h=50)
    para(s,MX,134,832,26,[(subtitle,INK60,False)],size=15,anchor=MSO_ANCHOR.TOP,ls=1.3)
    # задача — одразу після опису клієнта
    rect(s,MX,168,832,62,fill=MINT,radius=14)
    rect(s,MX+22,188,4,22,fill=BLUE)
    para(s,MX+40,168,772,62,[("Задача: ",INK,True),(task,INK60,False)],size=13.5,anchor=MSO_ANCHOR.MIDDLE,ls=1.28)
    gx,gy=MX,250; cw,ch=200,150
    for i,(n,l) in enumerate(stats4):
        x=gx+i*(cw+11)
        rect(s,x,gy,cw,ch,fill=TINT,radius=18)
        rect(s,x+22,gy+24,14,14,fill=GREEN)
        txt(s,x+22,gy+44,cw-40,52,n,size=36,bold=True,color=BLUE,anchor=MSO_ANCHOR.TOP)
        txt(s,x+22,gy+98,cw-36,50,l,size=11.5,color=INK60,ls=1.18)
    link(s,MX,438,320,"Детальніше про кейс  →",case_url,size=14,color=BLUE)
    link(s,512,438,320,"Читати всі кейси  →","https://site24.com.ua/tag/kejsi/",size=14,color=INK60,align=PP_ALIGN.RIGHT)
    footer(s,foot_num)

# ===================================================================== 01 ТИТУЛ
s=slide(TINT)
pic_h(s,os.path.join(IC,"site24-ink.png"),MX,58,32)
pill(s,MX,168,"ПРО КОМПАНІЮ",size=13)
heading(s,MX,206,820,[("Агенція ",INK,True),("digital-маркетингу",BLUE,True)],size=46,h=120)
para(s,MX,326,720,60,[("SEO та PPC просування, що приводить клієнтів з Google та AI",INK60,False)],size=20,anchor=MSO_ANCHOR.TOP)
# зелена дія-плашка (акцент 10%)
rect(s,MX,398,300,4,fill=GREEN)
txt(s,MX,504,460,20,[("site24.com.ua",INK60,False),("      ·      ",INK12,False),
    ("+38 (098) 738 77 08",INK60,True)],size=13,align=PP_ALIGN.LEFT)
txt(s,720,504,180,20,"2026",size=12,color=INK60,align=PP_ALIGN.RIGHT)

# ===================================================================== 02 ХТО МИ + STATS
s=slide(WHITE)
heading(s,MX,54,820,[("Хто ",INK,True),("ми",BLUE,True)],size=30)
para(s,MX,118,832,84,[("Site24 — агенція digital-маркетингу, що спеціалізується на просуванні сайтів у Google "
    "через SEO та PPC. Працюємо як ",INK,False),("партнер",INK,True),
    (", який допомагає малому та середньому бізнесу залучати клієнтів із пошуку — без зайвих технічних клопотів.",INK,False)],
    size=17,anchor=MSO_ANCHOR.TOP,ls=1.4)
stats=[("1000+","сайтів вивели в топ Google"),("550+","прибуткових рекламних кампаній запущено"),
       ("86%","клієнтів приходять до нас за рекомендацією"),("9","років успіху на ринку просування")]
gx,gy=MX,236; cw=200; ch=150
for i,(n,l) in enumerate(stats):
    x=gx+i*(cw+11)
    rect(s,x,gy,cw,ch,fill=TINT,radius=18)
    rect(s,x+22,gy+26,14,14,fill=GREEN)         # зелений прапорець-маркер
    txt(s,x+22,gy+46,cw-40,54,n,size=42,bold=True,color=BLUE,anchor=MSO_ANCHOR.TOP)
    txt(s,x+22,gy+100,cw-36,48,l,size=12.5,color=INK60,ls=1.18)
footer(s,"02")

# ===================================================================== 03 КОНКУРЕНТНІ ПЕРЕВАГИ
s=slide(WHITE)
heading(s,MX,50,860,[("Конкурентні ",INK,True),("переваги",BLUE,True)],size=30)
adv=[("Партнер, а не підрядник",
      "Не просто виконуємо ТЗ, а самі пропонуємо покращення й беремо на себе повний цикл просування: "
      "лінкбілдинг, крауд-маркетинг, UX, роботу з розробниками та аналітику. Готові «підставити плече» "
      "на кожному етапі й вести проєкт як власний."),
     ("Орієнтація на ROI",
      "Будуємо стратегію так, щоб SEO та PPC приносили максимальну окупність — не просто позиції в топі, "
      "а більше клієнтів і прибутку. Перші результати клієнт бачить уже за 2–3 місяці, і вони працюють "
      "у довгостроковій перспективі."),
     ("Прозорість",
      "Клієнт завжди знає, що ми робимо, за що платить і які результати отримуємо. Регулярна звітність "
      "зрозумілою мовою бізнесу — конкретні цифри та вигоди, без складних технічних термінів.")]
gx,gy=MX,120; cw,ch=270,252
for i,(t,d) in enumerate(adv):
    x=gx+i*(cw+11)
    rect(s,x,gy,cw,ch,fill=MINT,radius=18)
    rect(s,x+20,gy+22,38,4,fill=BLUE)
    txt(s,x+20,gy+42,cw-38,52,t,size=18,bold=True,color=INK,ls=1.1)
    txt(s,x+20,gy+104,cw-38,ch-118,d,size=13,color=INK60,ls=1.35)
footer(s,"03")

# ===================================================================== 04 КЛІЄНТСЬКИЙ СЕРВІС
s=slide(WHITE)
heading(s,MX,50,890,[("Чому клієнтський сервіс ",INK,True),("Site24",BLUE,True),(" найкращий",INK,True)],size=27)
serv=[("Постійний контакт з менеджером",
       [("Ваш персональний project-менеджер завжди на зв'язку, володіє всією актуальною інформацією та "
         "оперативно відповідає на запити. ",INK60,False),
        ("90% звернень вирішуються в день надходження.",BLUE,True)]),
      ("Прозоре спілкування без «води»",
       [("Пояснюємо складні речі простою мовою, без технічного жаргону. Ви завжди розумієте, що і навіщо "
         "ми робимо у вашій рекламній кампанії.",INK60,False)]),
      ("Відкритість у бюджетах та звітах",
       [("Працюємо з відкритими бюджетами: ви чітко бачите, скільки йде на рекламу, а скільки — на нашу "
         "роботу. Усі дії задокументовані та прозорі.",INK60,False)])]
gx,gy=MX,116; cw,ch=270,232
for i,(t,d) in enumerate(serv):
    x=gx+i*(cw+11)
    rect(s,x,gy,cw,ch,fill=MINT,radius=18)
    rect(s,x+20,gy+22,38,4,fill=BLUE)
    txt(s,x+20,gy+40,cw-38,50,t,size=17,bold=True,color=INK,ls=1.12)
    txt(s,x+20,gy+96,cw-38,ch-110,d,size=12.5,color=INK60,ls=1.32)
rect(s,MX,372,832,50,fill=BLUE,radius=14)
tag=para(s,MX,372,832,50,[("Ваш успіх — ",WHITE,False),("наша репутація.",GREENM,False)],
     size=17,anchor=MSO_ANCHOR.MIDDLE)
tag.text_frame.paragraphs[0].alignment=PP_ALIGN.CENTER
footer(s,"04")

# ===================================================================== 05 СЕРТИФІКАТИ
s=slide(WHITE)
heading(s,MX,54,860,[("Сертифікати",BLUE,True),(" та статуси",INK,True)],size=30)
para(s,MX,116,832,30,[("Підтверджена експертиза й офіційні партнерські статуси провідних платформ.",INK60,False)],
     size=15,anchor=MSO_ANCHOR.TOP)
certs=[("google-partner.png","Офіційний партнерський статус Google Ads — доступ до бета-функцій, навчання та підтримки платформи."),
       ("clutch.jpg","Верифікований профіль із реальними відгуками клієнтів на міжнародному рейтингу агенцій."),
       ("horoshop.png","Сертифікований підрядник для e-commerce на одній із провідних платформ України.")]
gx,gy=MX,176; cw,ch=270,210
for i,(logo,d) in enumerate(certs):
    x=gx+i*(cw+11)
    rect(s,x,gy,cw,ch,fill=WHITE,line=LINE,radius=18)
    rect(s,x,gy,cw,6,fill=BLUE)
    rect(s,x+24,gy+28,cw-48,84,fill=TINT,radius=12)
    fit_pic(s,os.path.join(BADGE,logo),x+40,gy+42,cw-80,56)
    txt(s,x+24,gy+128,cw-48,ch-146,d,size=13,color=INK60,ls=1.32)
footer(s,"05")

# ===================================================================== 06 ПОСЛУГИ ТА ІНСТРУМЕНТИ
s=slide(WHITE)
heading(s,MX,50,860,[("Послуги ",INK,True),("та технології",BLUE,True)],size=30)
svc=[("SEO-просування сайтів","Комплексне виведення сайту в топ Google за комерційними запитами"),
     ("PPC / контекстна реклама","Google Ads під ROI: пошук, банери, ремаркетинг, Shopping"),
     ("Просування в AI","Видимість у ChatGPT та AI Overview (AEO) — новий канал трафіку"),
     ("SERM — репутація в пошуку","Керування тим, що бачить клієнт про бренд у видачі"),
     ("Просування закордонних сайтів","SEO та PPC для інших країн, мов і ринків")]
gx,gy=MX,110; cw,ch=832,56
for i,(t,d) in enumerate(svc):
    y=gy+i*(ch+9)
    rect(s,gx,y,cw,ch,fill=MINT,radius=12)
    txt(s,gx+20,y,34,ch,[("→",BLUE,True)],size=18,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,gx+52,y,300,ch,t,size=16.5,bold=True,color=INK,anchor=MSO_ANCHOR.MIDDLE,ls=1.05)
    txt(s,gx+360,y,cw-380,ch,d,size=13.5,color=INK60,anchor=MSO_ANCHOR.MIDDLE,ls=1.1)
txt(s,MX,438,140,20,"Працюємо з:",size=13,bold=True,color=BLUE)
txt(s,MX+108,438,724,20,"Google Analytics 4  ·  Search Console  ·  Google Ads  ·  Ahrefs  ·  Serpstat  ·  Screaming Frog",
    size=13.5,color=INK60)
footer(s,"06")

# ===================================================================== 07 ГАЛУЗЕВИЙ ДОСВІД
s=slide(WHITE)
heading(s,MX,54,880,[("Галузевий ",INK,True),("досвід",BLUE,True)],size=30)
para(s,MX,116,832,30,[("Розуміємо специфіку різних ніш — від e-commerce до медицини — на ринку України та за кордоном.",INK60,False)],
     size=15,anchor=MSO_ANCHOR.TOP)
niches=[("E-commerce","ecommerce.png"),("Медицина","medicine.png"),("Виробництво","manufacturing.png"),
        ("Сервісні бізнеси","services.png"),("Нерухомість","realestate.png"),("IT та освіта","edtech.png")]
gx,gy=MX,172; cw,ch=270,74
for i,(t,ic) in enumerate(niches):
    x=gx+(i%3)*(cw+11); y=gy+(i//3)*(ch+14)
    rect(s,x,y,cw,ch,fill=TINT,radius=14)
    rect(s,x+16,y+(ch-46)/2,46,46,fill=WHITE,radius=11)
    fit_pic(s,os.path.join(NICHE,ic),x+16+9,y+(ch-46)/2+9,28,28)
    txt(s,x+78,y,cw-92,ch,t,size=17,bold=True,color=INK,anchor=MSO_ANCHOR.MIDDLE)
# банер міжнародного досвіду
rect(s,MX,400,832,64,fill=BLUE,radius=16)
para(s,MX+28,400,776,64,[("Локальний та міжнародний досвід — ",WHITE,False),
    ("працюємо і з українськими, і з закордонними проєктами: різні країни, мови та ринки.",GREENM,False)],size=15,anchor=MSO_ANCHOR.MIDDLE)
footer(s,"07")

# ===================================================================== 08 КЛІЄНТИ
s=slide(WHITE)
heading(s,MX,54,820,[("Нам ",INK,True),("довіряють",BLUE,True)],size=30)
para(s,MX,110,832,26,[("Понад 1000 бізнесів обрали Site24 для просування — від локальних магазинів до брендів національного рівня.",INK60,False)],
     size=15,anchor=MSO_ANCHOR.TOP)
lw,lh=imsize(os.path.join(AST,"clients-logos.jpeg")); iw=832; ih=iw*lh/lw
if ih>316: ih=316; iw=ih*lw/lh
s.shapes.add_picture(os.path.join(AST,"clients-logos.jpeg"),Pt(MX+(832-iw)/2),Pt(158),width=Pt(iw))
footer(s,"08")

# ===================================================================== 09 КЕЙС SEO
case_slide(
    "SEO-КЕЙС",
    [("Кейс: ",INK,True),("музичний інтернет-магазин",BLUE,True)],
    "Спеціалізований онлайн-магазин гітар та музичних інструментів на платформі Horoshop.",
    [("×24","зростання SEO-трафіку: 667 → 17 205 візитів/міс"),
     ("1922","запити в ТОП-3 Google (48% ядра)"),
     ("3705","запитів у ТОП-10 (92% ядра)"),
     ("18 міс","тривалість проєкту")],
    "вивести пріоритетні категорії гітар у ТОП-3 Google з мінімальним бюджетом на посилання — старт на початку повномасштабного вторгнення.",
    "https://site24.com.ua/seo-prodvizhenie-internet-magazina-muzykalnyh-instrumentov/",
    "09")

# ===================================================================== 10 КЕЙС PPC
case_slide(
    "PPC-КЕЙС",
    [("Кейс: ",INK,True),("e-commerce — дохід ×2,5 за рік",BLUE,True)],
    "The Tea (thetea.ua) — інтернет-магазин китайського чаю та аксесуарів для чайних церемоній.",
    [("×2,5","зростання доходу за рік"),
     ("×2,3","транзакції: 940 → 2143 /міс"),
     ("+200%","ROAS: 333% → 523%"),
     ("30+","нових Google Ads кампаній")],
    "масштабувати контекстну рекламу магазину чаю The Tea та збільшити кількість замовлень і дохід.",
    "https://site24.com.ua/ppc-dlya-e-commerce-zrostannya-dohodu-v-25-razi-za-odin-rik/",
    "10")

# ===================================================================== 11 КОНТАКТИ
s=slide(TINT)
pic_h(s,os.path.join(IC,"site24-ink.png"),MX,58,30)
heading(s,MX,150,760,[("Обговоримо ",INK,True),("ваш проєкт",BLUE,True),("?",INK,True)],size=38,h=60)
para(s,MX,222,700,30,[("Проведемо безкоштовну консультацію та покажемо точки зростання для вашого сайту.",INK60,False)],
     size=17,anchor=MSO_ANCHOR.TOP)
def lab(x,y,t): txt(s,x,y,360,18,t,size=12.5,bold=True,color=BLUE)
def val(x,y,parts,sz=18): txt(s,x,y,420,26,parts,size=sz,color=INK)
lab(MX,296,"КЕРІВНИК АГЕНЦІЇ")
val(MX,318,"Катерина Золотарьова")
lab(MX,372,"ТЕЛЕФОН")
val(MX,394,"+38 (098) 738 77 08")
lab(500,372,"САЙТ")
val(500,394,[("site24.com.ua",BLUE,True)])
# зелена дія-кнопка (акцент) — веде в Telegram власниці
TG="https://t.me/katerinazolotaryova"
btn=rect(s,500,296,300,52,fill=GREEN,radius=12)
btn.click_action.hyperlink.address=TG
btxt=txt(s,500,296,268,52,"Замовити консультацію",size=16,bold=True,color=INK,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
btxt.text_frame.paragraphs[0].runs[0].hyperlink.address=TG
txt(s,750,312,26,22,"→",size=18,bold=True,color=INK,anchor=MSO_ANCHOR.MIDDLE)
txt(s,60,504,160,20,"11",size=12,color=INK60,align=PP_ALIGN.LEFT)

out=os.path.join(BASE,"2026-07-15-EN-site24-company-short.pptx")
prs.save(out)
print("SAVED",out,"| slides",len(prs.slides._sldIdLst))
